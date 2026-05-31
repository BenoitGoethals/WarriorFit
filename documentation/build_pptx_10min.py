"""Build WarriorFit_Presentatie_10min_NL.pptx — 15 slides, 10-minute deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Palette — military olive / amber theme
# ---------------------------------------------------------------------------
C_BG        = RGBColor(0x1E, 0x2D, 0x1E)   # very dark olive (background)
C_PANEL     = RGBColor(0x2A, 0x40, 0x2A)   # panel bg
C_AMBER     = RGBColor(0xC8, 0xA0, 0x20)   # amber / gold
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_KHAKI     = RGBColor(0xB0, 0xC0, 0x9A)   # khaki subtitle
C_TBLHDR    = RGBColor(0x3A, 0x5A, 0x3A)   # table header row
C_TBLALT    = RGBColor(0x24, 0x38, 0x24)   # table alt row
C_RED       = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN     = RGBColor(0x28, 0xB4, 0x63)
C_BLUE      = RGBColor(0x21, 0x6A, 0xAB)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ---------------------------------------------------------------------------
# Low-level XML helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor | None = None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def txb(slide, text, x, y, w, h,
        font_size=20, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text-box and return the shape."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def txb_lines(slide, lines, x, y, w, h,
              font_size=18, bold=False, color=C_WHITE,
              align=PP_ALIGN.LEFT, spacing_after=6, line_color=None):
    """Add multiple lines to one textbox; line_color can be a list."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = "Calibri"
        if isinstance(line_color, list):
            run.font.color.rgb = line_color[i] if i < len(line_color) else color
        else:
            run.font.color.rgb = line_color if line_color else color
    return box


def title_bar(slide, title, subtitle=None):
    """Amber top bar with title + optional subtitle."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=C_AMBER)
    txb(slide, title,
        Inches(0.35), Inches(0.1), Inches(12.5), Inches(0.8),
        font_size=32, bold=True, color=C_BG)
    if subtitle:
        txb(slide, subtitle,
            Inches(0.35), Inches(0.8), Inches(12.5), Inches(0.4),
            font_size=15, bold=False, color=C_BG)


def add_table(slide, headers, rows, x, y, w, h,
              col_widths=None, hdr_size=14, row_size=13):
    """Add a styled table."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    # column widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    # header
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TBLHDR
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(hdr_size)
        run.font.color.rgb = C_AMBER
        run.font.name = "Calibri"
        _set_cell_margins(cell)
    # data rows
    for ri, row in enumerate(rows):
        bg = C_TBLALT if ri % 2 == 0 else C_PANEL
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(row_size)
            run.font.color.rgb = C_WHITE
            run.font.name = "Calibri"
            _set_cell_margins(cell)
    return tbl


def _set_cell_margins(cell, top=Pt(3), bot=Pt(3), left=Pt(5), right=Pt(5)):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for attr, val in [("marT", int(top)), ("marB", int(bot)),
                      ("marL", int(left)), ("marR", int(right))]:
        tcPr.set(attr, str(val))


def bullet(slide, items, x, y, w, h, font_size=19, color=C_WHITE,
           bullet_char="▶", bullet_color=C_AMBER, spacing_after=8):
    """Bullet list with colored bullet character."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing_after)
        # bullet run
        rb = p.add_run()
        rb.text = bullet_char + "  "
        rb.font.size = Pt(font_size)
        rb.font.color.rgb = bullet_color
        rb.font.name = "Calibri"
        # text run
        rt = p.add_run()
        if isinstance(item, tuple):
            rt.text = item[0]
            rt.font.color.rgb = item[1]
        else:
            rt.text = item
            rt.font.color.rgb = color
        rt.font.size = Pt(font_size)
        rt.font.name = "Calibri"
    return box


def kpi_strip(slide, kpis, y=Inches(1.4)):
    """Row of KPI boxes: list of (value, label) tuples."""
    n = len(kpis)
    box_w = Inches(13.33 / n)
    for i, (val, lbl) in enumerate(kpis):
        x = Inches(i * 13.33 / n)
        panel = add_rect(slide, x, y, box_w, Inches(1.2), fill=C_PANEL)
        txb(slide, val, x, y + Inches(0.08), box_w, Inches(0.65),
            font_size=30, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
        txb(slide, lbl, x, y + Inches(0.65), box_w, Inches(0.45),
            font_size=13, bold=False, color=C_KHAKI, align=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 1 — Title
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)

# centre panel
add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(4.8), fill=C_PANEL)

# logo text / emblem placeholder
txb(slide, "⚔", Inches(5.8), Inches(1.6), Inches(1.7), Inches(1.2),
    font_size=60, bold=False, color=C_AMBER, align=PP_ALIGN.CENTER)

txb(slide, "WarriorFit",
    Inches(1), Inches(2.7), Inches(11.33), Inches(1.2),
    font_size=54, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

txb(slide, "Digitalisering van militaire fysieke fitheidstesten",
    Inches(1), Inches(3.75), Inches(11.33), Inches(0.7),
    font_size=24, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)

txb(slide, "Benoit Goethals  |  mei 2026",
    Inches(1), Inches(4.45), Inches(11.33), Inches(0.5),
    font_size=17, color=C_KHAKI, align=PP_ALIGN.CENTER)

# stats bar at bottom
kpi_strip(slide, [
    ("1.012", "commits"),
    ("25.800", "lijnen Python"),
    ("231", "story points"),
    ("19", "sprints"),
    ("80", "user stories"),
], y=Inches(5.8))


# ===========================================================================
# SLIDE 2 — Probleemstelling
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Probleemstelling",
          "Manuele, papieren registratie van militaire fitheidstesten")

col_x1 = Inches(0.4)
col_x2 = Inches(7.0)
col_w  = Inches(6.2)

# Left column — problem list
add_rect(slide, col_x1, Inches(1.4), col_w, Inches(5.7), fill=C_PANEL)
txb(slide, "Huidige situatie", col_x1 + Inches(0.2), Inches(1.5),
    col_w - Inches(0.3), Inches(0.5),
    font_size=17, bold=True, color=C_AMBER)
bullet(slide, [
    "Papieren fiches per test — geen centraal overzicht",
    "5 testtypes, elk met eigen scoringsregels",
    "Manuele HR-synchronisatie (foutgevoelig, traag)",
    "Geen historisch inzicht per militair",
    "Geen AVG-compliance (gezondheidsdata)",
    "Geen rol-gebaseerde toegangscontrole",
    "Geen automatische herinnering of planning",
], col_x1 + Inches(0.2), Inches(2.05),
   col_w - Inches(0.3), Inches(4.8),
   font_size=18, spacing_after=10)

# Right column — impact
add_rect(slide, col_x2, Inches(1.4), col_w, Inches(5.7), fill=C_PANEL)
txb(slide, "Impact", col_x2 + Inches(0.2), Inches(1.5),
    col_w - Inches(0.3), Inches(0.5),
    font_size=17, bold=True, color=C_AMBER)
bullet(slide, [
    ("Vertraging in HR-rapportage", C_RED),
    ("Dataverlies door papieren archief", C_RED),
    ("Geen audit trail — inbreuk op regelgeving", C_RED),
    ("PTI's verliezen uren aan manueel werk", C_RED),
    ("Commandanten missen realtimezicht op eenheidsstatus", C_RED),
], col_x2 + Inches(0.2), Inches(2.05),
   col_w - Inches(0.3), Inches(3.5),
   font_size=18, spacing_after=12)

txb(slide, "→  WarriorFit digitaliseert het volledige proces",
    col_x2 + Inches(0.2), Inches(5.8),
    col_w - Inches(0.3), Inches(0.8),
    font_size=17, bold=True, color=C_GREEN)


# ===========================================================================
# SLIDE 3 — Oplossing & Ecosysteem
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Oplossing — WarriorFit Ecosysteem",
          "Intranetwebplatform — geen publieke blootstelling")

kpi_strip(slide, [
    ("24", "pagina's"),
    ("6", "rollen"),
    ("5", "testtypes"),
    ("40+", "DI-singletons"),
    ("100 %", "async I/O"),
], y=Inches(1.35))

# Component boxes row
components = [
    ("Shiny for Python", "Reactieve web-UI\n24 pagina's\nRBAC-beschermd"),
    ("PostgreSQL\n+ asyncpg", "Async ORM\nSQLAlchemy 2.0\nAlembic-migraties"),
    ("Broker (MOM)", "Transactionele outbox\nExp. backoff retry\nDead-letter queue"),
    ("FastAPI", "MOM REST-endpoint\nX-API-Key auth\nCORS-vergrendeling"),
    ("HR Simulator\n+ Mailpit", "Dev/test-omgeving\nHR-integratie mock\nE-mailpreview"),
]
bw = Inches(2.45)
bh = Inches(2.7)
by = Inches(2.85)
for i, (title, body) in enumerate(components):
    bx = Inches(0.25 + i * 2.6)
    add_rect(slide, bx, by, bw, bh, fill=C_PANEL)
    txb(slide, title, bx + Inches(0.1), by + Inches(0.1),
        bw - Inches(0.2), Inches(0.7),
        font_size=15, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)
    txb_lines(slide, body.split("\n"),
              bx + Inches(0.1), by + Inches(0.75),
              bw - Inches(0.2), bh - Inches(0.9),
              font_size=13, color=C_KHAKI, align=PP_ALIGN.CENTER)

txb(slide, "Deployment: intranet — geen publiek internet · Docker · GitHub Actions CI/CD",
    Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.4),
    font_size=13, italic=True, color=C_KHAKI, align=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 4 — Architectuur
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Architectuur — Lagen & DI Container",
          "Unidirectionele afhankelijkheidsstroom")

# Architecture layers
layers = [
    ("UI — Shiny pagina's", "24 modules · @inject · Provide[Container.xxx]", C_AMBER, Inches(1.3)),
    ("Controllers",         "PhefController, CrossController, ... (5 controllers)", RGBColor(0x8E, 0x44, 0xAD), Inches(2.1)),
    ("Services",           "ServiceTest, MilitaryService, GdprService, ... (12+)", C_BLUE, Inches(2.9)),
    ("Repositories",       "FitnessTestRepository, UserRepository, ... (8 repos)", RGBColor(0xE6, 0x7E, 0x22), Inches(3.7)),
    ("ORM Models",         "PhefTest, FunctionalTest, March, ... — polymorf", RGBColor(0x27, 0xAE, 0x60), Inches(4.5)),
    ("PostgreSQL",         "Async · asyncpg · TLS verify-full (productie)", C_KHAKI, Inches(5.3)),
]
for name, desc, color, y in layers:
    add_rect(slide, Inches(1.0), y, Inches(7.5), Inches(0.65), fill=C_PANEL)
    txb(slide, name, Inches(1.15), y + Inches(0.05),
        Inches(3.2), Inches(0.55),
        font_size=16, bold=True, color=color)
    txb(slide, desc, Inches(4.4), y + Inches(0.12),
        Inches(4.0), Inches(0.55),
        font_size=13, color=C_KHAKI)

# Arrows between layers
for iy in [1.95, 2.75, 3.55, 4.35, 5.15]:
    txb(slide, "↓", Inches(4.15), Inches(iy), Inches(0.5), Inches(0.3),
        font_size=14, color=C_AMBER, align=PP_ALIGN.CENTER)

# DI panel on right
add_rect(slide, Inches(9.0), Inches(1.35), Inches(4.0), Inches(5.8), fill=C_PANEL)
txb(slide, "DI Container", Inches(9.1), Inches(1.4), Inches(3.8), Inches(0.5),
    font_size=16, bold=True, color=C_AMBER)
txb_lines(slide, [
    "providers.Singleton(…)",
    "",
    "ApplicationConfig",
    "8 repositories",
    "12+ services",
    "5 controllers",
    "MailService · BEMILService",
    "NotifyMail · Broker",
    "",
    "→ 40+ singletons",
    "→ 1× aangemaakt bij opstart",
    "→ gedeeld over volledige sessie",
], Inches(9.2), Inches(1.95), Inches(3.6), Inches(5.0),
   font_size=13, color=C_KHAKI)

txb(slide, "→ 40+", Inches(9.2), Inches(5.85), Inches(1.5), Inches(0.4),
    font_size=20, bold=True, color=C_AMBER)


# ===========================================================================
# SLIDE 5 — Rollen & RBAC
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Rollen & Toegangscontrole (RBAC)",
          "6 rollen · PageSpec met allowed_roles per module")

headers = ["Rol", "Omschrijving", "Toegang"]
rows = [
    ("ADMIN",    "Systeembeheerder", "Volledige toegang + beheer"),
    ("PTI",      "Fitheidsinstructeur", "Testresultaten invoeren & bekijken"),
    ("APTI",     "Adjunct PTI", "Beperkt invoer, volledige lezing"),
    ("PLANNER",  "Planningsofficial", "Sessies & fitnessruimte plannen"),
    ("GUEST",    "Gast / bezoeker", "Alleen eigen profiel (read-only)"),
    ("USER",     "Militair (serviceman)", "Eigen fiche, privacypagina"),
]
add_table(slide, headers, rows,
          Inches(0.4), Inches(1.4), Inches(7.8), Inches(4.8),
          col_widths=[Inches(1.6), Inches(2.4), Inches(3.8)],
          hdr_size=15, row_size=14)

# Right panel
add_rect(slide, Inches(8.5), Inches(1.4), Inches(4.5), Inches(5.8), fill=C_PANEL)
txb(slide, "Implementatie", Inches(8.7), Inches(1.5), Inches(4.2), Inches(0.5),
    font_size=16, bold=True, color=C_AMBER)
bullet(slide, [
    "PageSpec(allowed_roles=[…]) per pagina",
    "Sessiecontrole bij elke navigatie",
    "10-min inactiviteitstime-out",
    "Argon2id wachtwoordhashing",
    "Rate-limiting op login",
    "Audit trail: alle CRUD-acties",
    "UserStore scoped per Shiny-sessie",
    "IDOR-bescherming (_assert_can_modify)",
], Inches(8.7), Inches(2.1), Inches(4.1), Inches(4.8),
   font_size=14, spacing_after=7)


# ===========================================================================
# SLIDE 6 — Testtypes
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "5 Testtypes — volledig gedigitaliseerd",
          "ORM: polymorf FitnessTest-basisklasse → subtypes")

tests = [
    ("PHEF", "Fysieke Fitheid",
     "2400 m looptijd\nZijbrug (links & rechts)\nLeeftijds-/gendercorrectie\nScores: A / B / C / D / F"),
    ("Functioneel", "Functionele Test",
     "Functionele fitheid\nKomplexe oefeningen\nRij scores 1-5\nGeschikt voor hersteltraject"),
    ("Gevechtstest", "Paraat (Paratrooper)",
     "Gevechtsgerichte test\nParatrooper-specifiek\nHarde drempelwaarden\nBinaire uitkomst"),
    ("Mars", "Marstest",
     "Timed march met uitrusting\nAfstand + gewicht\nEenheidsprestatie\nHistoriek per militair"),
    ("Cross", "Crossloop",
     "Tijdmetingen per loper\nChronos XML-import\nStatistieken (8 tabbladen)\nPodium & trends"),
]
bw = Inches(2.4)
bh = Inches(4.8)
by = Inches(1.4)
colors = [C_AMBER, C_GREEN, C_RED, C_BLUE, RGBColor(0x8E, 0x44, 0xAD)]
for i, (code, full, desc) in enumerate(tests):
    bx = Inches(0.25 + i * 2.6)
    add_rect(slide, bx, by, bw, Inches(0.55), fill=colors[i])
    txb(slide, code, bx, by, bw, Inches(0.55),
        font_size=20, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, by + Inches(0.55), bw, bh - Inches(0.55), fill=C_PANEL)
    txb(slide, full, bx + Inches(0.1), by + Inches(0.65), bw - Inches(0.2), Inches(0.5),
        font_size=13, bold=True, color=colors[i])
    txb_lines(slide, desc.split("\n"),
              bx + Inches(0.1), by + Inches(1.2), bw - Inches(0.2), Inches(3.3),
              font_size=13, color=C_KHAKI)

txb(slide, "Gemeenschappelijk: ORM-polymorfisme · async repository · audit trail · HR-broker",
    Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.4),
    font_size=13, italic=True, color=C_KHAKI, align=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 7 — PHEF Dataflow & Broker
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Dataflow — PHEF-invoer & Broker (MOM)",
          "Transactionele outbox · exponentieel backoff · dead-letter")

# Synchronous path
txb(slide, "Synchroon pad  (gebruiker wacht)",
    Inches(0.4), Inches(1.4), Inches(6.4), Inches(0.45),
    font_size=15, bold=True, color=C_AMBER)
sync_steps = [
    "Browser (PTI klikt 'Add')",
    "→  PhefPage (Shiny reactive)",
    "→  PhefController.add_phef()",
    "→  ServiceTest.add_phef_test()",
    "→  FitnessTestRepository.add()",
    "→  PostgreSQL — INSERT",
    "←  OK + grid refresh",
]
add_rect(slide, Inches(0.4), Inches(1.9), Inches(6.4), Inches(3.6), fill=C_PANEL)
txb_lines(slide, sync_steps, Inches(0.6), Inches(2.0),
          Inches(6.0), Inches(3.4),
          font_size=16, color=C_WHITE, spacing_after=9)

# Broker path
txb(slide, "Asynchroon pad  (broker — gebruiker wacht NIET)",
    Inches(7.0), Inches(1.4), Inches(6.0), Inches(0.45),
    font_size=15, bold=True, color=C_GREEN)
broker_steps = [
    "Service → Broker.send_message(test)",
    "→  asyncio.Queue  (in-memory buffer)",
    "→  _process_cycle()  elke 5 s",
    "→  hr_messages table  (duurzaam)",
    "→  POST naar HR-systeem",
    "✓  success → rij verwijderd",
    "✗  failure → retry (exp. backoff 5 s→10 min)",
    "✗✗ max_attempts → dead-letter = true",
]
add_rect(slide, Inches(7.0), Inches(1.9), Inches(6.0), Inches(3.8), fill=C_PANEL)
txb_lines(slide, broker_steps, Inches(7.2), Inches(2.0),
          Inches(5.7), Inches(3.6),
          font_size=15, spacing_after=8,
          line_color=[C_WHITE, C_WHITE, C_WHITE, C_WHITE, C_WHITE,
                      C_GREEN, C_RED, C_RED])

# Bottom note
add_rect(slide, Inches(0.4), Inches(5.9), Inches(12.6), Inches(0.8), fill=C_TBLHDR)
txb(slide, ("UI geeft onmiddellijk feedback — HR-synchronisatie gebeurt op de achtergrond. "
            "Broker overleeft uitval van het HR-systeem van willekeurige duur."),
    Inches(0.6), Inches(5.95), Inches(12.2), Inches(0.65),
    font_size=14, italic=True, color=C_KHAKI)


# ===========================================================================
# SLIDE 8 — Agile / Scrum metrics
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Agile Methodologie — Scrum metrics",
          "Sep 2025 – mei 2026 · 1 developer · 2-weekse sprints")

kpi_strip(slide, [
    ("20", "epics"),
    ("80", "user stories"),
    ("231", "story points"),
    ("19", "sprints"),
    ("12,2", "gem. velocity SP"),
    ("720 h", "geschatte inspanning"),
], y=Inches(1.35))

# Sprint table (sample)
headers = ["Sprint", "Periode", "SP", "Cumul.", "Focus"]
rows = [
    ("S1",  "sep 01–14",  "10", "10",  "Projectopzet, auth, login"),
    ("S2",  "sep 15–28",  "13", "23",  "Gebruikersbeheer, BEMIL"),
    ("S5",  "okt 27–nov 09", "18", "74",  "4 epics — hoogste dichtheid (160 commits)"),
    ("S7",  "nov 24–dec 07", "16", "110", "Cross sessies & statistieken"),
    ("S12", "feb 02–15",  "12", "162", "DI-container refactor"),
    ("S17", "apr 13–26",  "25", "221", "GDPR + cross-stats redesign"),
    ("S19", "mei 11–24",  "10", "231", "Beveiliging, DPIA, finale"),
]
add_table(slide, headers, rows,
          Inches(0.4), Inches(2.8), Inches(8.8), Inches(4.0),
          col_widths=[Inches(0.9), Inches(1.5), Inches(0.7), Inches(0.85), Inches(4.85)],
          hdr_size=14, row_size=12)

# Epic top 5
add_rect(slide, Inches(9.4), Inches(2.8), Inches(3.7), Inches(4.0), fill=C_PANEL)
txb(slide, "Top 5 epics (55 % backlog)", Inches(9.5), Inches(2.9), Inches(3.5), Inches(0.4),
    font_size=13, bold=True, color=C_AMBER)
top_epics = [
    ("E8  Cross & Statistieken", "28 SP"),
    ("E20 GDPR / Privacy",       "25 SP"),
    ("E1  Gebruikersbeheer",     "18 SP"),
    ("E3  PHEF-invoer",          "18 SP"),
    ("E2  Sessieplanning",       "15 SP"),
]
for i, (ep, sp) in enumerate(top_epics):
    y0 = Inches(3.4 + i * 0.62)
    add_rect(slide, Inches(9.5), y0, Inches(3.5), Inches(0.55), fill=C_TBLALT)
    txb(slide, ep,  Inches(9.6), y0 + Inches(0.05), Inches(2.4), Inches(0.45), font_size=12, color=C_WHITE)
    txb(slide, sp,  Inches(12.0), y0 + Inches(0.05), Inches(0.9), Inches(0.45), font_size=13, bold=True, color=C_AMBER, align=PP_ALIGN.RIGHT)


# ===========================================================================
# SLIDE 9 — Kwaliteit & Testing
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Kwaliteit & Testing",
          "pytest-asyncio · mypy strict · ruff · pre-commit · GitHub Actions")

col_w = Inches(6.0)
# Left — testing
add_rect(slide, Inches(0.4), Inches(1.4), col_w, Inches(5.8), fill=C_PANEL)
txb(slide, "Testpiramide", Inches(0.6), Inches(1.5), Inches(5.6), Inches(0.45),
    font_size=16, bold=True, color=C_AMBER)
bullet(slide, [
    "pytest + pytest-asyncio voor async tests",
    "Unittest.mock: AsyncMock, MagicMock",
    "Constructor-injectie maakt mocking triviaal",
    "container.override() voor integratietests",
    "reset_override() — geen statuslek tussen tests",
    "Singleton._instances clearzen in fixtures",
], Inches(0.6), Inches(2.0), Inches(5.6), Inches(3.8),
   font_size=16, spacing_after=10)
txb(slide, "Voorbeeld:",
    Inches(0.6), Inches(5.4), Inches(5.6), Inches(0.35),
    font_size=13, bold=True, color=C_AMBER)
txb(slide, "PhefController(service=mock_svc, mil_service=mock_mil)",
    Inches(0.6), Inches(5.75), Inches(5.6), Inches(0.55),
    font_size=12, color=RGBColor(0x8B, 0xD1, 0xFF), italic=True)

# Right — tooling
add_rect(slide, Inches(6.8), Inches(1.4), col_w, Inches(5.8), fill=C_PANEL)
txb(slide, "Tooling & CI/CD", Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.45),
    font_size=16, bold=True, color=C_AMBER)
tools = [
    ("mypy",         "Strenge typechecking — strict mode"),
    ("ruff",         "Linting: E/F/W/I · regellengte 100"),
    ("black",        "Codeformattering"),
    ("pre-commit",   "Auto-update version.yaml bij commit"),
    ("GitHub Actions", "CI: lint + mypy + tests bij elke push"),
    ("Docker",       "Productie-image · geen secrets in image"),
    ("Alembic",      "DB-migraties versiebeheerd"),
]
for i, (tool, desc) in enumerate(tools):
    y0 = Inches(2.05 + i * 0.72)
    txb(slide, tool, Inches(7.0), y0, Inches(1.8), Inches(0.62),
        font_size=14, bold=True, color=C_AMBER)
    txb(slide, desc, Inches(8.85), y0, Inches(3.7), Inches(0.62),
        font_size=14, color=C_WHITE)


# ===========================================================================
# SLIDE 10 — Beveiliging (OWASP)
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Beveiliging — OWASP Top 10 evaluatie",
          "Productie-intranet · defense-in-depth aanpak")

headers = ["OWASP risico", "Maatregel in WarriorFit", "Status"]
rows = [
    ("A01 — Broken Access Control",  "RBAC per PageSpec · IDOR-fix · sessie-scoped UserStore", "✓ gedekt"),
    ("A02 — Cryptographic Failures", "Argon2id hashing · PostgreSQL TLS verify-full",           "✓ gedekt"),
    ("A03 — Injection",              "SQLAlchemy parameterized queries · geen raw SQL",           "✓ gedekt"),
    ("A05 — Security Misconfiguration", "CORS-vergrendeling · X-API-Key op MOM · strict CSP",   "✓ gedekt"),
    ("A07 — Auth Failures",          "Rate-limiting login · 10-min timeout · audit trail",       "✓ gedekt"),
    ("A09 — Logging & Monitoring",   "Volledige CRUD-audit log · broker-health monitoring",      "✓ gedekt"),
    ("A10 — Server-Side Request Forgery", "HR-URL allowlist (config) — HTTPS open punt",        "⚠ open"),
]
add_table(slide, headers, rows,
          Inches(0.4), Inches(1.4), Inches(12.6), Inches(5.4),
          col_widths=[Inches(3.0), Inches(7.2), Inches(2.4)],
          hdr_size=14, row_size=13)

# colour status cells
tbl_shape = slide.shapes[-1]
tbl = tbl_shape.table
for ri in range(1, len(rows) + 1):
    cell = tbl.cell(ri, 2)
    val  = rows[ri - 1][2]
    cell.fill.solid()
    cell.fill.fore_color.rgb = C_GREEN if "✓" in val else C_RED
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = C_BG


# ===========================================================================
# SLIDE 11 — AVG / GDPR Compliance
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "AVG / GDPR Compliance",
          "DPIA uitgevoerd · Art. 35 · gezondheidsdata Art. 9")

col_w = Inches(6.0)

add_rect(slide, Inches(0.4), Inches(1.4), col_w, Inches(5.7), fill=C_PANEL)
txb(slide, "Rechtsgrond & scope", Inches(0.6), Inches(1.5), col_w - Inches(0.3), Inches(0.45),
    font_size=15, bold=True, color=C_AMBER)
bullet(slide, [
    "Art. 6(1)(c) — wettelijke verplichting (militaire opleiding)",
    "Art. 6(1)(e) — taak van algemeen belang",
    "Art. 9(2)(b) — gezondheidsdata in arbeidscontext",
    "DPIA uitgevoerd (Art. 35) — versie 1.1, 30 apr 2026",
    "Verwerkingsverantwoordelijke: Belgische Defensie",
    "Uitsluitend intranet — geen publieke blootstelling",
], Inches(0.6), Inches(2.0), col_w - Inches(0.3), Inches(4.8),
   font_size=15, spacing_after=10)

add_rect(slide, Inches(6.8), Inches(1.4), col_w, Inches(5.7), fill=C_PANEL)
txb(slide, "Geïmplementeerde rechten", Inches(7.0), Inches(1.5), col_w - Inches(0.3), Inches(0.45),
    font_size=15, bold=True, color=C_AMBER)
rights = [
    ("Art. 15 — Inzage",     "Volledig exporteerbaar via zelfservice Privacy-pagina"),
    ("Art. 17 — Verwijdering", "GdprService.erase_user() · FK CASCADE"),
    ("Art. 20 — Portabiliteit", "JSON-export van eigen fiche"),
    ("Art. 7  — Toestemming", "user_consents tabel · intrekbaar"),
    ("RetentionService",     "Automatische verwijdering na configureerbare termijn"),
    ("Audit trail",          "Alle CRUD-acties gelogd (Art. 32-vereiste)"),
]
for i, (art, desc) in enumerate(rights):
    y0 = Inches(2.05 + i * 0.82)
    txb(slide, art, Inches(7.0), y0, Inches(2.0), Inches(0.75),
        font_size=13, bold=True, color=C_GREEN)
    txb(slide, desc, Inches(9.05), y0, Inches(3.55), Inches(0.75),
        font_size=13, color=C_WHITE)


# ===========================================================================
# SLIDE 12 — Technologiestapel
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Technologiestapel",
          "Python 3.13+ · moderne async stack · volledig open source")

tech_groups = [
    ("Frontend / UI", C_AMBER, [
        "Shiny for Python — reactieve server-side UI",
        "Plotly — interactieve grafieken",
        "Rajdhani + JetBrains Mono — militaire typografie",
        "CSS-thema: olive / khaki / amber",
    ]),
    ("Backend / Data", C_GREEN, [
        "PostgreSQL + asyncpg",
        "SQLAlchemy 2.0 Async ORM",
        "Alembic — DB-migratiebeheer",
        "dependency-injector — DI container",
    ]),
    ("Beveiliging", C_RED, [
        "Argon2id — wachtwoordhashing",
        "RBAC per PageSpec (6 rollen)",
        "TLS verify-full (productie)",
        "X-API-Key + CORS-vergrendeling",
    ]),
    ("Infrastructuur / Dev", C_BLUE, [
        "Docker + GitHub Actions CI/CD",
        "FastAPI — MOM REST endpoint",
        "Mailpit — e-mailpreview (dev)",
        "HR Simulator — integratiemock",
        "uv — dependency management",
        "ruff · mypy · black · pre-commit",
    ]),
]

cols = 2
gw = Inches(6.2)
gh = Inches(2.6)
for idx, (group, color, items) in enumerate(tech_groups):
    col = idx % cols
    row = idx // cols
    gx = Inches(0.4 + col * 6.5)
    gy = Inches(1.5 + row * 2.85)
    add_rect(slide, gx, gy, gw, gh, fill=C_PANEL)
    add_rect(slide, gx, gy, gw, Inches(0.45), fill=color)
    txb(slide, group, gx + Inches(0.15), gy + Inches(0.05),
        gw - Inches(0.2), Inches(0.38),
        font_size=14, bold=True, color=C_BG)
    txb_lines(slide, items,
              gx + Inches(0.15), gy + Inches(0.5),
              gw - Inches(0.2), gh - Inches(0.6),
              font_size=13, color=C_KHAKI, spacing_after=4)


# ===========================================================================
# SLIDE 13 — Ontwikkelingstraject (fasen)
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Ontwikkelingstraject — 5 fasen",
          "Sep 2025 – mei 2026 · 8,5 maanden · 1 developer")

phases = [
    ("Fase 1", "sep–okt 2025", "Prototype",
     "Basisauth, login-modal\nGebruikersbeheer (E1)\nRolmodel opgezet",
     RGBColor(0x5D, 0x6D, 0x7E)),
    ("Fase 2", "nov 2025–jan 2026", "Feature Growth",
     "Alle 5 testtypes\nBEMIL-integratie\nCross-sessies (E8)",
     RGBColor(0x27, 0x6F, 0xAB)),
    ("Fase 3", "feb 2026", "Architectuur",
     "DI-container refactor\n40+ singletons\nAsync repository-patroon",
     RGBColor(0x8E, 0x44, 0xAD)),
    ("Fase 4", "mrt–apr 2026", "Hardening",
     "Broker transactionele outbox\nGDPR-compliance (E20)\nCross-statistieken (8 tabs)",
     RGBColor(0xE6, 0x7E, 0x22)),
    ("Fase 5", "apr–mei 2026", "Security",
     "OWASP Top 10\nDPIA uitvoering\nPostgreSQL TLS\nFinale oplevering",
     C_AMBER),
]
bw = Inches(2.4)
bh = Inches(4.7)
by = Inches(1.5)
for i, (fase, period, title, desc, color) in enumerate(phases):
    bx = Inches(0.25 + i * 2.6)
    add_rect(slide, bx, by, bw, Inches(0.55), fill=color)
    txb(slide, fase, bx, by, bw, Inches(0.55),
        font_size=16, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, bx, by + Inches(0.55), bw, bh - Inches(0.55), fill=C_PANEL)
    txb(slide, period, bx + Inches(0.1), by + Inches(0.65), bw - Inches(0.2), Inches(0.4),
        font_size=11, italic=True, color=C_KHAKI)
    txb(slide, title, bx + Inches(0.1), by + Inches(1.1), bw - Inches(0.2), Inches(0.5),
        font_size=14, bold=True, color=color)
    txb_lines(slide, desc.split("\n"),
              bx + Inches(0.1), by + Inches(1.65), bw - Inches(0.2), bh - Inches(1.9),
              font_size=13, color=C_KHAKI, spacing_after=6)

# timeline arrow
add_rect(slide, Inches(0.4), Inches(6.48), Inches(12.5), Inches(0.12), fill=C_AMBER)
txb(slide, "→ tijdlijn", Inches(12.3), Inches(6.3), Inches(1.0), Inches(0.3),
    font_size=11, color=C_AMBER, align=PP_ALIGN.RIGHT)


# ===========================================================================
# SLIDE 14 — Routekaart
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)
title_bar(slide, "Routekaart — verdere ontwikkeling",
          "Open punten & toekomstige iteraties")

items = [
    (C_RED,    "Hoog", "SSO / SAML-integratie",
     "Eenmalige aanmelding via Defensie-identiteitsprovider"),
    (C_RED,    "Hoog", "HR-URL HTTPS allowlist (OWASP A10)",
     "Whitelist HR-endpoints + HTTPS-verificatie op uitgaand verkeer"),
    (C_AMBER,  "Gemiddeld", "Serviceman login — wachtwoordverificatie",
     "SSO of eigen credentials afwachten; nu tijdelijk overgeslagen (GitHub-issue)"),
    (C_AMBER,  "Gemiddeld", "Responsief ontwerp (mobiel)",
     "Tablet-/smartphoneweergave voor PTI's op terrein"),
    (C_BLUE,   "Laag", "Realtimenotificaties",
     "WebSocket push voor sessiewijzigingen en broker-alerts"),
    (C_BLUE,   "Laag", "Multi-eenheid ondersteuning",
     "Uitbreiding naar meerdere regimenten / compagnieën"),
    (RGBColor(0x27, 0xAE, 0x60), "Laag", "Uitgebreid rapportagepakket",
     "PDF-exports per test, eenheid en tijdsperiode"),
]
for i, (color, prio, title, desc) in enumerate(items):
    y0 = Inches(1.5 + i * 0.74)
    add_rect(slide, Inches(0.4), y0, Inches(1.3), Inches(0.62), fill=color)
    txb(slide, prio, Inches(0.4), y0 + Inches(0.07), Inches(1.3), Inches(0.5),
        font_size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.75), y0, Inches(11.0), Inches(0.62), fill=C_PANEL)
    txb(slide, title, Inches(1.9), y0 + Inches(0.07), Inches(3.2), Inches(0.5),
        font_size=14, bold=True, color=C_WHITE)
    txb(slide, desc, Inches(5.15), y0 + Inches(0.1), Inches(7.5), Inches(0.48),
        font_size=13, color=C_KHAKI, italic=True)


# ===========================================================================
# SLIDE 15 — Afsluiting / Vragen
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_BG)

add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_BG)
add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.2), fill=C_PANEL)

txb(slide, "⚔", Inches(5.8), Inches(1.3), Inches(1.7), Inches(1.0),
    font_size=52, color=C_AMBER, align=PP_ALIGN.CENTER)

txb(slide, "Bedankt!",
    Inches(1), Inches(2.2), Inches(11.33), Inches(1.1),
    font_size=52, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

txb(slide, "Vragen?",
    Inches(1), Inches(3.1), Inches(11.33), Inches(0.7),
    font_size=28, bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)

txb(slide, "Benoit Goethals  ·  benoit.goethals@gmail.com",
    Inches(1), Inches(4.0), Inches(11.33), Inches(0.5),
    font_size=17, color=C_KHAKI, align=PP_ALIGN.CENTER)

txb(slide, "github.com/BenoitGoethals/WarriorFit",
    Inches(1), Inches(4.5), Inches(11.33), Inches(0.45),
    font_size=15, italic=True, color=C_KHAKI, align=PP_ALIGN.CENTER)

kpi_strip(slide, [
    ("20", "epics"),
    ("80", "stories"),
    ("231 SP", "geleverd"),
    ("1.012", "commits"),
    ("25.800", "lijnen Python"),
    ("8,5 mnd", "solo-ontwikkeling"),
], y=Inches(5.8))

# ---------------------------------------------------------------------------
out = "/Users/benoit/PycharmProjects/WarriorFit/documentation/WarriorFit_Presentatie_10min_NL.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
